import os
import sqlite3
from dataclasses import dataclass
from typing import Dict, List, Tuple

from flask import Flask, render_template, request, url_for

import yfinance as yf
import pandas as pd
import matplotlib
matplotlib.use("Agg")  # 不需要顯示視窗（避免部署時出錯）
import matplotlib.pyplot as plt
import numpy as np

import time
from datetime import datetime

import glob

from pptx import Presentation
from pptx.util import Inches, Pt

# =========================
# 核心回測資料結構
# =========================
@dataclass
class BacktestResult:
    symbol: str
    start_year: int
    end_year: int
    monthly_invest: float
    months: int

    strategy_a_final_asset: float
    strategy_a_return: float

    strategy_b_final_asset: float
    strategy_b_return: float

    bear_years: list

    portfolio_summary: list    
    portfolio_details: dict

    # 曲線：如果你後面想做更多圖/表可用
    curve: Dict[str, List[float]]


# =========================
# (1) SQLite：啟動時建 DB
# =========================
def init_db(db_path: str = "history.db") -> None:
    conn = sqlite3.connect(db_path)
    cur = conn.cursor()
    cur.execute(
        """
        CREATE TABLE IF NOT EXISTS history (
            id INTEGER PRIMARY KEY AUTOINCREMENT,
            symbol TEXT NOT NULL,
            monthly_invest REAL NOT NULL,
            start_year INTEGER NOT NULL,
            end_year INTEGER NOT NULL,

            strategy_a_final_asset REAL NOT NULL,
            strategy_a_return REAL NOT NULL,

            strategy_b_final_asset REAL NOT NULL,
            strategy_b_return REAL NOT NULL,

            created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
        );
        """
    )
    conn.commit()
    conn.close()


def insert_history(
    symbol: str,
    monthly_invest: float,
    start_year: int,
    end_year: int,
    a_final: float,
    a_return: float,
    b_final: float,
    b_return: float,
    db_path: str = "history.db",
) -> None:
    conn = sqlite3.connect(db_path)
    cur = conn.cursor()
    cur.execute(
        """
        INSERT INTO history (
            symbol, monthly_invest, start_year, end_year,
            strategy_a_final_asset, strategy_a_return,
            strategy_b_final_asset, strategy_b_return
        ) VALUES (?, ?, ?, ?, ?, ?, ?, ?);
        """,
        (symbol, monthly_invest, start_year, end_year, a_final, a_return, b_final, b_return),
    )
    conn.commit()
    conn.close()

def get_recent_history(limit: int = 5, db_path: str = "history.db"):
    conn = sqlite3.connect(db_path)
    cur = conn.cursor()
    cur.execute(
        """
        SELECT symbol, monthly_invest, start_year, end_year,
               strategy_a_final_asset, strategy_a_return,
               strategy_b_final_asset, strategy_b_return,
               created_at
        FROM history
        ORDER BY created_at DESC
        LIMIT ?;
        """,
        (limit,),
    )
    rows = cur.fetchall()
    conn.close()

    # rows: (symbol, monthly_invest, start_year, end_year, a_final, a_return, b_final, b_return, created_at)
    history_list = []
    for r in rows:
        history_list.append({
            "symbol": r[0],
            "monthly_invest": r[1],
            "start_year": r[2],
            "end_year": r[3],
            "strategy_a_final_asset": r[4],
            "strategy_a_return": r[5],
            "strategy_b_final_asset": r[6],
            "strategy_b_return": r[7],
            "created_at": r[8],
        })
    return history_list



# =========================
# (2) 回測核心函式：Adj Close + 每月最後一個交易日
# =========================
def backtest_dca_vs_lump_sum_adjclose(
    symbols: List[str],
    monthly_invest: float,
    invest_week: int,
    invest_day_type: str
) -> Tuple[BacktestResult, "matplotlib.figure.Figure"]:
    
    if monthly_invest <= 0:
        raise ValueError("每月投資金額 monthly_invest 必須 > 0")

    start_year = 2018
    end_year = datetime.now().year
    start_date = "2018-01-01"
    end_date = datetime.now().strftime("%Y-%m-%d")

    all_symbols = symbols + ["^TWII"]
    df = yf.download(all_symbols, start=start_date, end=end_date, progress=False, auto_adjust=False)

    if df is None or df.empty:
        raise ValueError(f"查無資料：({start_date} ~ {end_date})")

    adj_close = df["Adj Close"]
    daily_close = adj_close.ffill().bfill()
    
    selected_dates = []
    for (y, m), group in daily_close.groupby([daily_close.index.year, daily_close.index.month]):
        week_series = group.index.day.map(lambda d: min((d - 1) // 7 + 1, 4))
        week_data = group[week_series == invest_week]
        
        if not week_data.empty:
            selected_dates.append(week_data.index[0] if invest_day_type == "first" else week_data.index[-1])
        else:
            selected_dates.append(group.index[0] if invest_day_type == "first" else group.index[-1])
                
    monthly_prices = daily_close.loc[selected_dates]
    dates = monthly_prices.index.to_list()
    months = len(dates)
    total_invested = months * monthly_invest

    # ---- 策略 A & B 資產計算 ----
    monthly_per_stock = monthly_invest / len(symbols)
    a_shares = {sym: 0.0 for sym in symbols}
    a_assets_curve = []
    
    # 用來記錄每檔股票、每個月的交易明細
    portfolio_details = {sym: [] for sym in symbols}

    remaining_cash = total_invested

    # 1. 歷經每個月，計算定期定額買入股數與當下市值
    for date_idx, row in monthly_prices.iterrows():
        monthly_val = 0.0
        date_str = date_idx.strftime("%Y-%m-%d")
        remaining_cash -= monthly_invest
        for sym in symbols:
            price = row[sym]
            shares_bought = monthly_per_stock / price
            a_shares[sym] += shares_bought
            monthly_val += a_shares[sym] * price
            
            # 記錄這筆買進的詳細資訊
            portfolio_details[sym].append({
                "date": date_str,
                "buy_price": float(price),
                "shares": float(shares_bought),
                "cost": float(monthly_per_stock)
            })
            
        a_assets_curve.append(monthly_val + remaining_cash)

    strategy_a_final_asset = float(a_assets_curve[-1])
    strategy_a_return = (strategy_a_final_asset - total_invested) / total_invested

    # 2. 結算目前每檔股票的總損益，並回推每個月買進批次的現值
    portfolio_summary = []
    final_prices = {sym: float(monthly_prices.iloc[-1][sym]) for sym in symbols}
    
    for sym in symbols:
        tot_shares = a_shares[sym]
        tot_cost = monthly_per_stock * months
        curr_price = final_prices[sym]
        
        # 💎 計算均價成本 (加上防呆機制)
        avg_cost = tot_cost / tot_shares if tot_shares > 0 else 0
        
        mkt_val = tot_shares * curr_price
        pl = mkt_val - tot_cost
        pl_pct = (pl / tot_cost) * 100 if tot_cost > 0 else 0
        
        # 建立總表
        portfolio_summary.append({
            "symbol": sym,
            "shares": tot_shares,
            "lots": tot_shares / 1000.0, 
            "cost": tot_cost,
            "avg_cost": avg_cost,        # 寫入均價成本
            "current_price": curr_price,
            "market_value": mkt_val,
            "pl": pl,
            "pl_pct": pl_pct
        })
        
        # 幫每個月的明細加上「放到今天的現值與損益」
        for tx in portfolio_details[sym]:
            tx_curr_val = tx["shares"] * curr_price
            tx_pl = tx_curr_val - tx["cost"]
            tx_pl_pct = (tx_pl / tx["cost"]) * 100 if tx["cost"] > 0 else 0
            tx["current_value"] = tx_curr_val
            tx["pl"] = tx_pl
            tx["pl_pct"] = tx_pl_pct

    b_prices = monthly_prices["^TWII"].to_list()
    first_price = b_prices[0]
    b_shares = total_invested / first_price
    b_assets_curve = [b_shares * p for p in b_prices]

    strategy_b_final_asset = float(b_assets_curve[-1])
    strategy_b_return = (strategy_b_final_asset - total_invested) / total_invested

    # --- 💎 新增：計算年度報酬率與空頭年分析 ---
    yearly_returns = {}
    years_list = sorted(list(set(d.year for d in dates)))
    
    for y in years_list:
        mask = [d.year == y for d in dates]
        if sum(mask) < 2: continue # 資料不足跳過
        first_idx = mask.index(True)
        last_idx = len(mask) - 1 - mask[::-1].index(True)
        
        # 該年大盤報酬
        twii_start, twii_end = b_prices[first_idx], b_prices[last_idx]
        twii_ret = (twii_end - twii_start) / twii_start * 100
        
        # 該年組合報酬 (等權重)
        port_ret = 0
        for sym in symbols:
            sym_prices = [row[sym] for idx, row in monthly_prices.iterrows()]
            s_start, s_end = sym_prices[first_idx], sym_prices[last_idx]
            port_ret += ((s_end - s_start) / s_start * 100) / len(symbols)
            
        yearly_returns[y] = {'portfolio': port_ret, 'twii': twii_ret}
        
    bar_years = list(yearly_returns.keys())
    bar_port = [yearly_returns[y]['portfolio'] for y in bar_years]
    bar_twii = [yearly_returns[y]['twii'] for y in bar_years]
    
    # 篩選出大盤跌幅 < 0 的空頭年份
    bear_years_data = [{'year': y, 'twii': yearly_returns[y]['twii'], 'port': yearly_returns[y]['portfolio']} 
                       for y in bar_years if yearly_returns[y]['twii'] < 0]

    # --- 💎 修改：畫出雙圖表 (左邊折線、右邊柱狀) ---
    plt.rcParams['font.sans-serif'] = ['Microsoft JhengHei']  
    plt.rcParams['axes.unicode_minus'] = False

    fig, (ax1, ax2) = plt.subplots(2, 1, figsize=(12, 10))
    display_name = ", ".join(symbols)
    
    # 左圖：折線圖
    ax1.plot(dates, a_assets_curve, label=f"策略 A：組合", linewidth=2)
    ax1.plot(dates, b_assets_curve, label="策略 B：大盤 ^TWII", linewidth=2)
    ax1.set_title("資產成長曲線 (2018 至今)")
    ax1.set_xlabel("時間")
    ax1.set_ylabel("總資產")
    ax1.legend()
    ax1.grid(True, alpha=0.3)
    
    # 右圖：柱狀圖
    x = np.arange(len(bar_years))
    width = 0.35
    ax2.bar(x - width/2, bar_port, width, label='策略 A (組合)')
    ax2.bar(x + width/2, bar_twii, width, label='策略 B (大盤)')
    ax2.set_xticks(x)
    ax2.set_xticklabels(bar_years)
    ax2.set_title("年度績效比較 (%)")
    ax2.legend()
    ax2.grid(True, alpha=0.3)
    
    plt.tight_layout()

    result = BacktestResult(
        symbol=display_name, 
        start_year=start_year,
        end_year=end_year,
        monthly_invest=monthly_invest,
        months=months,
        strategy_a_final_asset=strategy_a_final_asset,
        strategy_a_return=float(strategy_a_return),
        strategy_b_final_asset=strategy_b_final_asset,
        strategy_b_return=float(strategy_b_return),
        bear_years=bear_years_data,  # 傳入空頭年資料
        portfolio_summary=portfolio_summary,  
        portfolio_details=portfolio_details,
        curve={"dates": [], "a_assets": [], "b_assets": []}
    )
    return result, fig

def get_stock_info(symbol: str) -> dict:
    """
    取得股票基本面資訊（用 yf.Ticker(symbol).info）
    回傳格式：
    {
      "shortName": str 或 "無資料",
      "trailingPE": str 或 "無資料",
      "dividendYield": str 或 "無資料",  # 百分比字串
      "marketCap": str 或 "無資料"
    }
    """
    try:
        ticker = yf.Ticker(symbol)
        info = ticker.info or {}
    except Exception:
        info = {}

    def get_or_na(key):
        v = info.get(key, None)
        return v if v not in (None, "") else "無資料"

    short_name = get_or_na("shortName")

    trailing_pe = info.get("trailingPE", None)
    if trailing_pe is None or trailing_pe == "":
        trailing_pe_str = "無資料"
    else:
        # trailingPE 可能是 float，轉成整數或保留 2 位都可以
        try:
            trailing_pe_str = f"{float(trailing_pe):,.2f}"
        except Exception:
            trailing_pe_str = str(trailing_pe)

    dividend_yield = info.get("dividendYield", None)  # 通常是 0.xx
    if dividend_yield is None or dividend_yield == "":
        dividend_yield_str = "無資料"
    else:
        try:
            dividend_yield_str = f"{float(dividend_yield) * 100:.2f}%"
        except Exception:
            dividend_yield_str = "無資料"

    market_cap = info.get("marketCap", None)  # 可能是整數（美元）
    if market_cap is None or market_cap == "":
        market_cap_str = "無資料"
    else:
        try:
            # 以「美金」做簡單千/百萬/十億顯示（更好看）
            mc = float(market_cap)
            if mc >= 1e12:
                market_cap_str = f"{mc/1e12:.2f}T"
            elif mc >= 1e9:
                market_cap_str = f"{mc/1e9:.2f}B"
            elif mc >= 1e6:
                market_cap_str = f"{mc/1e6:.2f}M"
            else:
                market_cap_str = f"{mc:,.0f}"
        except Exception:
            market_cap_str = "無資料"

    return {
        "shortName": short_name,
        "trailingPE": trailing_pe_str,
        "dividendYield": dividend_yield_str,
        "marketCap": market_cap_str,
    }

def create_ppt_report(result: BacktestResult, img_path: str, ppt_path: str):
    prs = Presentation()
    
    # ---------------- 投影片 1：標題頁 ----------------
    slide_title = prs.slides.add_slide(prs.slide_layouts[0])
    slide_title.shapes.title.text = "投資組合回測與大盤對決報告"
    slide_title.placeholders[1].text = f"分析標的：{result.symbol}\n回測期間：{result.start_year} - {result.end_year}\n每月投入：{result.monthly_invest:,.0f} 元"

    # ---------------- 投影片 2：圖表與數據 ----------------
    slide_chart = prs.slides.add_slide(prs.slide_layouts[5])  # 只有標題的版面
    slide_chart.shapes.title.text = "績效對決與資產成長曲線"
    
    # 插入報酬率文字
    txBox = slide_chart.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.8))
    tf = txBox.text_frame
    tf.text = f"策略 A (組合) 總報酬率: {result.strategy_a_return*100:.2f}%  |  策略 B (大盤) 總報酬率: {result.strategy_b_return*100:.2f}%"
    
    # 💎 修正：只限制高度 (height=Inches(5.5))，讓寬度自動等比例縮放，保證圖片不變形也不會超出邊界！
    slide_chart.shapes.add_picture(img_path, Inches(0.5), Inches(1.8), height=Inches(5.5))

    # ---------------- 投影片 3：投資組合明細表 (全新加入) ----------------
    slide_table = prs.slides.add_slide(prs.slide_layouts[5])
    slide_table.shapes.title.text = "投資組合總持股明細"
    
    rows = len(result.portfolio_summary) + 2
    cols = 6
    table_shape = slide_table.shapes.add_table(rows, cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.8))
    table = table_shape.table
    
    headers = ["標的代號", "累積股數", "均價成本", "總成本", "目前市值", "未實現損益"]
    for i, h in enumerate(headers):
        table.cell(0, i).text = h
        
    for r_idx, item in enumerate(result.portfolio_summary):
        row = r_idx + 1
        table.cell(row, 0).text = item['symbol']
        table.cell(row, 1).text = f"{item['shares']:,.0f} 股"
        table.cell(row, 2).text = f"{item['avg_cost']:,.2f}"
        table.cell(row, 3).text = f"{item['cost']:,.0f}"
        table.cell(row, 4).text = f"{item['market_value']:,.0f}"
        table.cell(row, 5).text = f"{item['pl']:,.0f} ({item['pl_pct']:.2f}%)"
        
    total_cost = sum(item['cost'] for item in result.portfolio_summary)
    total_mkt = sum(item['market_value'] for item in result.portfolio_summary)
    total_pl = sum(item['pl'] for item in result.portfolio_summary)
    total_pl_pct = (total_pl / total_cost * 100) if total_cost > 0 else 0
    
    last_row = rows - 1
    table.cell(last_row, 0).text = "總計"
    table.cell(last_row, 1).text = "-"
    table.cell(last_row, 2).text = "-"
    table.cell(last_row, 3).text = f"{total_cost:,.0f}"
    table.cell(last_row, 4).text = f"{total_mkt:,.0f}"
    table.cell(last_row, 5).text = f"{total_pl:,.0f} ({total_pl_pct:.2f}%)"

    # ---------------- 投影片 4：智能診斷分析結論 (搬移網頁邏輯) ----------------
    slide_analysis = prs.slides.add_slide(prs.slide_layouts[1]) # 標題與內容版面
    slide_analysis.shapes.title.text = "💡 綜合結論與投資啟示"
    
    tf_analysis = slide_analysis.placeholders[1].text_frame
    tf_analysis.clear() # 清空預設文字
    
    a_asset = result.strategy_a_final_asset
    b_asset = result.strategy_b_final_asset
    a_pct = result.strategy_a_return * 100
    b_pct = result.strategy_b_return * 100
    bear_years = result.bear_years

    p1 = tf_analysis.add_paragraph()
    p1.font.size = Pt(18)
    p1.font.bold = True
    p1.text = "1. 整體績效評估："

    p2 = tf_analysis.add_paragraph()
    p2.font.size = Pt(16)
    
    if a_asset > b_asset:
        p2.text = "綜合來看，策略 A（投資組合）整體的投資績效優於策略 B（台灣大盤）。這表示您挑選的標的具備強大的成長動能，為您帶來了超越市場平均的超額報酬。"
    else:
        p2.text = "綜合來看，策略 B（台灣大盤）整體的投資績效優於策略 A（投資組合）。這顯示在本次回測期間，直接投資大盤能獲得最穩健且優渥的市場報酬。"
        
        lagging_stocks = [item for item in result.portfolio_summary if item['pl_pct'] < b_pct]
        p2_lag = tf_analysis.add_paragraph()
        p2_lag.font.size = Pt(16)
        if lagging_stocks:
            p2_lag.text = f"⚠️ 拖累組合績效的標的："
            for item in lagging_stocks:
                p_stock = tf_analysis.add_paragraph()
                p_stock.level = 1
                p_stock.font.size = Pt(14)
                p_stock.text = f"{item['symbol']}：總報酬率僅 {item['pl_pct']:.2f}%"
        else:
            p_stock = tf_analysis.add_paragraph()
            p_stock.level = 1
            p_stock.font.size = Pt(14)
            p_stock.text = "無特定落後標的（可能是因定期定額初期資金運用率較單筆投入低所致）"

    p3 = tf_analysis.add_paragraph()
    p3.font.size = Pt(18)
    p3.font.bold = True
    p3.text = "\n2. 多空頭市場特性分析："

    p4 = tf_analysis.add_paragraph()
    p4.font.size = Pt(16)

    if a_asset > b_asset:
        if bear_years:
            lose_count = sum(1 for item in bear_years if item['port'] < item['twii'])
            if lose_count > 0:
                p4.text = "原則上波動較大的標的，在「多頭年」時表現會比大盤更亮眼，但在「空頭年」遇到市場回檔時，往往也會跌得更深。"
                p5 = tf_analysis.add_paragraph()
                p5.font.size = Pt(16)
                p5.text = f"以本次回測的空頭年為例，您的投資組合有 {lose_count} 次跌幅深於大盤，印證了高波動的特性。但策略 A 採用了「定期定額」的機制，雖然空頭時帳面跌幅深，卻能趁機「在低檔累積大量便宜股數」，這是多頭反彈時爆發優異績效的關鍵！"
            else:
                p4.text = "進一步分析發現，您的投資組合不僅在整體績效上勝出，在遇到大盤下跌的「空頭年」時，表現也全數優於大盤！"
                p5 = tf_analysis.add_paragraph()
                p5.font.size = Pt(16)
                p5.text = "這顯示您所選擇的投資標的具有「進可攻、退可守」的防禦與成長雙重優勢，在任何市場環境下都能穩定勝出，是一個極為優秀的投資組合！"
        else:
            p4.text = "原則上波動較大的標的，在「多頭年」時表現會比大盤更亮眼，但在「空頭年」遇到市場回檔時，往往也會跌得更深。不過在此回測期間內皆為多頭市場，您的投資組合展現了極強的爆發力，順利擊敗大盤。"
    else:
        if bear_years:
            win_count = sum(1 for item in bear_years if item['port'] > item['twii'])
            if win_count > 0:
                p4.text = f"進一步分析發現，在遇到市場下跌的「空頭年」中，策略 A 有 {win_count} 次表現優於大盤。顯示您的投資組合在空頭時較為抗跌，具備良好的防禦屬性；但代價是在多頭年時表現不佳、缺乏爆發力，導致最終總結算時整體績效輸給了大盤。"
            else:
                p4.text = "進一步分析發現，策略 A 不論在多頭年或空頭年，表現都比大盤（策略 B）還要差。這顯示該投資組合的績效明顯落後於市場平均。建議您可以重新檢視持股，適度增加「成長型」的標的，以提升整體的資產增長潛力。"
        else:
            p4.text = "在此回測期間內皆為多頭市場，而策略 A 未能超越大盤，顯示該組合可能缺乏足夠的成長動能。建議可適度增加「成長型」標的。"

    # 存檔
    prs.save(ppt_path)

# =========================
# Flask App
# =========================
app = Flask(__name__)

# DB 初始化：程式啟動時建立/更新 history.db
init_db("history.db")

# static 圖片固定輸出檔名（依你的需求）
RESULT_IMG_PATH = os.path.join("static", "result.png")


@app.route("/", methods=["GET"])
def index():
    return render_template("index.html")


@app.route("/backtest", methods=["GET", "POST"])
def backtest():
    # 接收表單欄位（多檔股票字串）
    symbols_raw = request.values.get("symbols", "").strip()
    symbol_list = [s.strip() for s in symbols_raw.split(",") if s.strip()]
    
    monthly_invest_raw = request.values.get("monthly_invest")
    invest_week = int(request.values.get("invest_week", 1))
    invest_day_type = request.values.get("invest_day_type", "last")

    if monthly_invest_raw is None:
        return render_template("result.html", error="缺少回測參數：monthly_invest")

    monthly_invest = float(monthly_invest_raw)
    start_year = 2018
    end_year = datetime.now().year

    if not symbol_list:
        return render_template("result.html", error="請輸入至少一檔股票代號。")

    try:
        # 1. 執行核心計算
        result, fig = backtest_dca_vs_lump_sum_adjclose(
            symbols=symbol_list,
            monthly_invest=monthly_invest,
            invest_week=invest_week,
            invest_day_type=invest_day_type
        )

        # 2. 確保 static 資料夾存在並清空舊圖表
        os.makedirs("static", exist_ok=True)
        for old_img in glob.glob("static/result_*.png"):
            try:
                os.remove(old_img)
            except Exception:
                pass

        # 3. 建立帶有時間戳記的全新檔名
        ts = int(time.time())
        new_filename = f"result_{ts}.png"
        new_img_path = os.path.join("static", new_filename)

        # 4. 存圖 (加上 bbox_inches='tight' 確保雙圖表邊緣不會被切掉)
        fig.savefig(new_img_path, dpi=150, bbox_inches='tight')
        plt.close(fig)

        # ---------------- 新增：產生 PPT 報告 ----------------
        BASE_DIR = os.path.dirname(os.path.abspath(__file__))
        STATIC_DIR = os.path.join(BASE_DIR, "static")
        
        new_ppt_filename = f"report_{ts}.pptx"
        new_ppt_path = os.path.join(STATIC_DIR, new_ppt_filename)
        create_ppt_report(result, new_img_path, new_ppt_path)
        ppt_url = url_for("static", filename=new_ppt_filename)
        # --------------------------------------------------------

        # 5. 寫入資料庫
        insert_history(
            symbol=result.symbol,
            monthly_invest=result.monthly_invest,
            start_year=result.start_year,
            end_year=result.end_year,
            a_final=result.strategy_a_final_asset,
            a_return=result.strategy_a_return,
            b_final=result.strategy_b_final_asset,
            b_return=result.strategy_b_return,
            db_path="history.db",
        )

        # 6. 抓取基本面資訊
        stock_info = get_stock_info(symbol_list[0])

        # 7. 將「新的檔名」傳遞給網頁模板！(這就是剛才破圖的關鍵)
        img_url = url_for("static", filename=new_filename)

        return render_template(
            "result.html",
            symbol=result.symbol,
            start_year=result.start_year,
            end_year=result.end_year,
            monthly_invest=result.monthly_invest,
            strategy_a_final_asset=result.strategy_a_final_asset,
            strategy_a_return=result.strategy_a_return,
            strategy_b_final_asset=result.strategy_b_final_asset,
            strategy_b_return=result.strategy_b_return,
            bear_years=result.bear_years,
            portfolio_summary=result.portfolio_summary,  
            portfolio_details=result.portfolio_details,
            stock_info=stock_info,
            img_url=img_url,
            ppt_url=ppt_url
        )

    except Exception as e:
        return render_template("result.html", error=str(e))

if __name__ == "__main__":
    # 開發環境啟動
    app.run(debug=True)
